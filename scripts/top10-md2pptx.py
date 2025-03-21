from pptx import Presentation
from pptx.util import Inches

# Créer une nouvelle présentation PowerPoint
presentation = Presentation()

# Liste des vulnérabilités OWASP Top 10 pour LLM 2025
titles = [
    "LLM01: Injection de Prompt",
    "LLM02: Gestion de Sortie Non Sécurisée",
    "LLM03: Empoisonnement de la chaîne d'approvisionnement",
    "LLM04: Empoisonnement des données et des modèles",
    "LLM05: Gestion Inappropriée des Sorties",
    "LLM06: Autonomie Excessive",
    "LLM07: Fuite des Prompts Système",
    "LLM08: Faiblesses des Vecteurs et des Embeddings",
    "LLM09: Désinformation",
    "LLM10: Consommation Excessive"
]

# Descriptions brèves pour chaque vulnérabilité
descriptions = [
    "L'injection de prompt permet à un utilisateur malveillant de modifier le comportement d'un LLM.",
    "La gestion de sortie non sécurisée expose les applications à des risques XSS.",
    "L'empoisonnement de la chaîne d'approvisionnement affecte l'intégrité des données d'entraînement.",
    "L'empoisonnement des données et des modèles introduit des vulnérabilités dans les données d'entraînement.",
    "La gestion inappropriée des sorties expose les applications à divers risques de sécurité.",
    "L'autonomie excessive peut mener à des actions dommageables.",
    "La fuite des prompts système révèle des informations sensibles.",
    "Les faiblesses des vecteurs et des embeddings permettent d'injecter du contenu nuisible.",
    "La désinformation se propage via les contenus générés par les LLM.",
    "La consommation excessive entraîne des risques de dégradation du service."
]

# Créer un slide pour chaque titre
for i in range(len(titles)):
    # Ajouter un slide avec une mise en page par défaut
    slide_layout = presentation.slide_layouts[0]  # Utiliser la première mise en page (par défaut)
    slide = presentation.slides.add_slide(slide_layout)

    # Ajouter le titre au slide
    title_placeholder = slide.shapes.title
    title_placeholder.text = titles[i]

    # Ajouter une description brève
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(2)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = descriptions[i]

# Enregistrer la présentation
file_name = "OWASP_Top_10_LLM_2025.pptx"
presentation.save(file_name)

print(f"Fichier PPTX généré : {file_name}")
